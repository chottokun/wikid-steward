from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from wikid_steward.core.okf_converter import generate_okf_frontmatter
from wikid_steward.core.profiles import resolve_profile
from wikid_steward.core.slug import generate_slug


def create_sample_pdf(file_path: Path, text: str):
    """ReportLab でサンプルの PDF を動的生成"""
    c = canvas.Canvas(str(file_path), pagesize=letter)
    c.drawString(100, 750, text)
    c.save()


def create_sample_docx(file_path: Path, title: str):
    """python-docx でサンプルの Word 文書を動的生成"""
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph("This is a sample Word document for wikid-steward testing.")
    doc.save(str(file_path))


def create_sample_xlsx(file_path: Path):
    """openpyxl でサンプルの Excel 表計算文書を動的生成"""
    wb = Workbook()
    ws = wb.active
    ws.title = "BOM Data"
    ws.append(["ITEM", "Part Name", "Qty", "Spec"])
    ws.append(["01", "MicroController", 1, "ARM Cortex M4"])
    ws.append(["02", "Power Supply", 2, "12V 5A"])
    wb.save(str(file_path))


def create_sample_pptx(file_path: Path, title: str):
    """python-pptx でサンプルの PowerPoint スライドを動的生成"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    prs.save(str(file_path))


def test_variety_ingest_pipeline_all_formats(tmp_path: Path):
    """PDF, DOCX, XLSX, PPTX, サイドカーYAML付きファイルの多種多様な元データバリエーションのパイプライン検証"""

    base_raw = tmp_path / "_raw"
    tmp_path / "staging"
    tmp_path / "wiki"
    tmp_path / "raw_sources"

    # 1. 論文 PDF (papers/ フォルダ)
    pdf_paper = base_raw / "papers" / "llm_survey.pdf"
    pdf_paper.parent.mkdir(parents=True)
    create_sample_pdf(pdf_paper, "Large Language Models Survey Paper")

    # 2. CAD図面 PDF (drawings/ フォルダ)
    pdf_drawing = base_raw / "drawings" / "DWG-2026-X88.pdf"
    pdf_drawing.parent.mkdir(parents=True)
    create_sample_pdf(pdf_drawing, "ITEM 01 - Controller - Spec ARM - Qty 1")

    # 3. Word 文書 (documents/ フォルダ)
    docx_file = base_raw / "documents" / "specification.docx"
    docx_file.parent.mkdir(parents=True)
    create_sample_docx(docx_file, "System Requirements Specification")

    # 4. Excel 表計算 (sheets/ フォルダ)
    xlsx_file = base_raw / "sheets" / "parts_list.xlsx"
    xlsx_file.parent.mkdir(parents=True)
    create_sample_xlsx(xlsx_file)

    # 5. PowerPoint スライド (slides/ フォルダ)
    pptx_file = base_raw / "slides" / "architecture.pptx"
    pptx_file.parent.mkdir(parents=True)
    create_sample_pptx(pptx_file, "System Architecture Slides")

    # 6. サイドカー YAML 付きカスタム PDF (custom/ フォルダ)
    pdf_custom = base_raw / "custom" / "special_device.pdf"
    pdf_custom.parent.mkdir(parents=True)
    create_sample_pdf(pdf_custom, "Special Device Document")
    yaml_sidecar = pdf_custom.with_suffix(".yaml")
    yaml_sidecar.write_text(
        "profile: drawing\nimages_scale: 3.5\ncustom_metadata:\n  vendor: AcmeCorp\n",
        encoding="utf-8",
    )

    test_files = [
        (pdf_paper, "paper", "directory_policy"),
        (pdf_drawing, "drawing", "directory_policy"),
        (docx_file, "paper", "default"),  # デフォルト
        (xlsx_file, "spreadsheet", "directory_policy"),
        (pptx_file, "presentation", "directory_policy"),
        (pdf_custom, "drawing", "sidecar_yaml"),
    ]

    print("\n=== Testing Variety of Data Formats & Profiles ===")
    for raw_file, expected_profile, expected_source in test_files:
        profile, prof_source, custom_meta = resolve_profile(raw_file, base_raw)

        print(f"File: {raw_file.name} -> Profile: {profile.name} (Source: {prof_source})")

        assert profile.name == expected_profile
        assert prof_source == expected_source

        # OKF ヘッダー生成のテスト
        slug = generate_slug(str(raw_file.relative_to(base_raw).with_suffix("")))
        frontmatter = generate_okf_frontmatter(
            doc_id=slug,
            title=raw_file.stem,
            doc_type=profile.doc_type,
            source_path=str(Path("raw_sources") / raw_file.relative_to(base_raw)),
            profile_used=profile.name,
            profile_source=prof_source,
            custom_metadata=custom_meta,
        )

        assert f"profile_used: {expected_profile}" in frontmatter
        assert f"profile_source: {expected_source}" in frontmatter
